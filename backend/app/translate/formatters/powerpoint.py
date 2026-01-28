"""
PowerPoint 演示文稿格式处理器
支持 .pptx 文件的翻译，保持格式
"""
import uuid
import asyncio
from pathlib import Path
from typing import Callable, Optional
from pptx import Presentation

from . import BaseFormatter
from ...config import settings


class PowerPointFormatter(BaseFormatter):
    """
    PowerPoint 演示文稿处理器

    使用 python-pptx 库处理 PowerPoint 文件
    保持幻灯片布局、格式等
    """

    def translate(
        self,
        source_path: str,
        target_lang: str,
        ai_translator,
        progress_callback: Optional[Callable[[int, int], None]] = None
    ) -> str:
        """
        翻译 PowerPoint 演示文稿（同步包装器，调用异步方法）

        Args:
            source_path: 源文件路径
            target_lang: 目标语言
            ai_translator: AI 翻译器实例
            progress_callback: 进度回调函数

        Returns:
            str: 翻译结果文件路径
        """
        # 检查翻译器是否支持并发方法
        if hasattr(ai_translator, 'translate_batch_async_concurrent'):
            # 获取线程数配置
            thread_count = getattr(ai_translator, 'thread_count', 5)

            # 创建新的事件循环并运行异步方法（不关闭循环，避免 AsyncOpenAI 客户端引用错误）
            loop = asyncio.new_event_loop()
            asyncio.set_event_loop(loop)
            try:
                return loop.run_until_complete(
                    self.translate_async(source_path, target_lang, ai_translator, thread_count, progress_callback)
                )
            finally:
                # 不关闭循环，让 asyncio 自动管理
                asyncio.set_event_loop(None)
        else:
            # 使用原有同步逻辑
            # 加载演示文稿
            prs = Presentation(source_path)

            # 收集所有需要翻译的文本
            texts_to_translate = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts_to_translate.append({
                            'shape': shape,
                            'text': shape.text
                        })

            total_count = len(texts_to_translate)

            # 批量翻译
            batch_size = 10
            for i in range(0, total_count, batch_size):
                batch = texts_to_translate[i:i + batch_size]
                texts = [item['text'] for item in batch]

                # 翻译
                translated = ai_translator.translate_batch(texts, target_lang)

                # 更新文本框
                for j, item in enumerate(batch):
                    if j < len(translated):
                        # 保持格式，只更新文本
                        text_frame = item['shape'].text_frame
                        text_frame.text = translated[j]

                # 更新进度
                if progress_callback:
                    progress_callback(min(i + batch_size, total_count), total_count)

            # 保存结果
            result_path = self._generate_result_path(source_path)
            prs.save(result_path)

            return result_path

    async def translate_async(
        self,
        source_path: str,
        target_lang: str,
        ai_translator,
        thread_count: int = 5,
        progress_callback: Optional[Callable[[int, int], None]] = None
    ) -> str:
        """
        异步翻译 PowerPoint 演示文稿（支持并发）

        Args:
            source_path: 源文件路径
            target_lang: 目标语言
            ai_translator: AI 翻译器实例
            thread_count: 并发线程数
            progress_callback: 进度回调函数

        Returns:
            str: 翻译结果文件路径
        """
        # 加载演示文稿
        prs = Presentation(source_path)

        # 收集所有需要翻译的文本
        texts_to_translate = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts_to_translate.append({
                        'shape': shape,
                        'text': shape.text
                    })

        total_count = len(texts_to_translate)

        # 检查是否支持并发翻译
        if hasattr(ai_translator, 'translate_batch_async_concurrent'):
            # 使用并发翻译
            batch_size = 10  # 并发模式下可以增大批次
            for i in range(0, total_count, batch_size):
                batch = texts_to_translate[i:i + batch_size]
                texts = [item['text'] for item in batch]

                # 翻译
                translated = await ai_translator.translate_batch_async_concurrent(
                    texts=texts,
                    target_lang=target_lang,
                    max_concurrency=thread_count,
                    progress_callback=progress_callback
                )

                # 更新文本框
                for j, item in enumerate(batch):
                    if j < len(translated):
                        # 保持格式，只更新文本
                        text_frame = item['shape'].text_frame
                        text_frame.text = translated[j]
        else:
            # 降级到普通异步翻译
            for i, item in enumerate(texts_to_translate):
                translated = await ai_translator.translate_text_async(item['text'], target_lang)
                text_frame = item['shape'].text_frame
                text_frame.text = translated

                if progress_callback:
                    progress_callback(i + 1, total_count)

        # 保存结果
        result_path = self._generate_result_path(source_path)
        prs.save(result_path)

        return result_path

    def _generate_result_path(self, source_path: str) -> str:
        """生成结果文件路径"""
        source = Path(source_path)
        filename = f"{source.stem}_translated_{uuid.uuid4().hex[:8]}{source.suffix}"
        return str(settings.TRANSLATE_DIR / filename)
